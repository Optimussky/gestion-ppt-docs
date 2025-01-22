# pyinstaller --noconsole --onefile main.py
import flet as ft
from pptx import Presentation
import pandas as pd
import os
import subprocess

# Función para detectar etiquetas en un documento PowerPoint
def detect_tags(ppt_path):
    prs = Presentation(ppt_path)
    tags = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for word in shape.text.split():
                    if word.startswith('{{') and word.endswith('}}'):
                        tags.add(word.strip('{}'))
    return list(tags)

# Función para leer las columnas de un archivo Excel
def get_excel_columns(excel_path):
    df = pd.read_excel(excel_path)
    return df.columns.tolist(), df

# Función para reemplazar etiquetas preservando formato del documento PowerPoint
def replace_tags_in_ppt(ppt_path, output_path, data):
    prs = Presentation(ppt_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for key, value in data.items():
                    if isinstance(value, pd.Timestamp):
                        value = value.strftime('%d-%m-%Y')
                    if f'{{{{{key}}}}}' in shape.text:
                        shape.text = shape.text.replace(f'{{{{{key}}}}}', str(value))
    prs.save(output_path)

# Función principal para la aplicación
def main(page: ft.Page):
    page.title = "Automatización de PowerPoint-Excel   By Alberto Romero®™"
    page.scroll = "adaptive"

    # Contenedores para mostrar etiquetas, columnas y resultados
    ppt_tags = ft.Column()
    excel_columns = ft.Column()
    result_message = ft.Text()

    # Variables para almacenar rutas de archivos y datos
    ppt_file_path = None
    excel_file_path = None
    excel_data = None
    detected_tags = []

    def load_ppt_file(e):
        nonlocal ppt_file_path, detected_tags
        ppt_file_path = e.files[0].path
        detected_tags = detect_tags(ppt_file_path)
        ppt_tags.controls = [ft.Text(f"Etiqueta detectada: {tag}") for tag in detected_tags]
        page.update()

    def load_excel_file(e):
        nonlocal excel_file_path, excel_data
        excel_file_path = e.files[0].path
        columns, data = get_excel_columns(excel_file_path)
        excel_data = data
        excel_columns.controls = [ft.Text(f"Columna de Excel: {col}") for col in columns]
        page.update()

        # Validación de columnas
        validate_columns()

    def validate_columns():
        if ppt_file_path and excel_file_path:
            num_columns = len(excel_data.columns)
            num_tags = len(detected_tags)
            if num_columns != num_tags:
                result_message.value = f"Error: El número de columnas en Excel ({num_columns}) no coincide con el número de variables en PowerPoint ({num_tags})."
                result_message.color = "red"  # Cambiar el color a rojo
            else:
                result_message.value = f"Validación exitosa: {num_columns} columnas en Excel coinciden con {num_tags} variables en PowerPoint."
                result_message.color = "green"  # Cambiar el color a verde si es exitoso
            page.update()

    def process_and_generate(e):
        if not ppt_file_path or not excel_file_path:
            result_message.value = "Por favor, cargue un archivo PowerPoint y un archivo Excel primero."
            page.update()
            return

        output_dir = os.path.join(os.getcwd(), "output")
        os.makedirs(output_dir, exist_ok=True)

        for idx, row in excel_data.iterrows():
            output_path = os.path.join(output_dir, f"output_{idx + 1}.pptx")
            replace_tags_in_ppt(ppt_file_path, output_path, row.to_dict())

        result_message.value = f"Archivos generados exitosamente en la carpeta '{output_dir}'."
        result_message.color = "black"  # Restablecer el color a negro
        page.update()

        # Abrir la carpeta de salida
        subprocess.run(["explorer", os.path.abspath(output_dir)])

    def show_instructions(e):
        def close_instructions(_):
            instructions.open = False
            page.update()

        instructions = ft.AlertDialog(
            modal=True,
            title=ft.Text("Instrucciones de Uso"),
            content=ft.Text(
                """1.- Para usar esta aplicación es necesario tener 2 archivos base: 
                    - Plantilla en PowerPoint, Base de datos en Excel.\n\n"""
                """2.- El archivo de PowerPoint debe contener dentro de paréntesis {{}} dobles el nombre de cada una de las columnas que contenga el archivo con la base de Excel: 
                    {{NOMBRE_COLUMNA_EXCEL}}
                    
                , para sustituir los valores en la plantilla de PowerPoint.\n\n"""
                """3.- Los nombres de los campos de Excel deberán estar de preferencia en MAYÚSCULAS pero siempre sin espacios; en todo caso, se deben unir más de una palabra con guiones bajos:
                 LUGAR_NACIMIENTO."""
            ),
            actions=[
                ft.TextButton("Cerrar", on_click=close_instructions)
            ]
        )
        page.dialog = instructions
        page.dialog.open = True
        page.update()

    # Carga de archivos PowerPoint
    ppt_file_uploader = ft.FilePicker(on_result=load_ppt_file)
    page.overlay.append(ppt_file_uploader)

    # Carga de archivos Excel
    excel_file_uploader = ft.FilePicker(on_result=load_excel_file)
    page.overlay.append(excel_file_uploader)

    # Botones para cargar archivos, procesar datos y mostrar instrucciones
    upload_ppt_button = ft.ElevatedButton("Cargar Documento PowerPoint", on_click=lambda _: ppt_file_uploader.pick_files(
        allow_multiple=False, allowed_extensions=["pptx"]))

    upload_excel_button = ft.ElevatedButton("Cargar Archivo Excel", on_click=lambda _: excel_file_uploader.pick_files(
        allow_multiple=False, allowed_extensions=["xlsx"]))

    generate_button = ft.ElevatedButton("Generar Archivos", on_click=process_and_generate)
    instructions_button = ft.ElevatedButton("Instrucciones", on_click=show_instructions)

    # Layout de la aplicación
    page.add(
        ft.Text("Automatización de Documentos PowerPoint", style="headlineMedium"),
        ft.Row([upload_ppt_button, upload_excel_button, instructions_button], alignment="center"),
        ft.Divider(),
        ft.Text("Etiquetas detectadas en el documento PowerPoint:"),
        ppt_tags,
        ft.Divider(),
        ft.Text("Columnas detectadas en el archivo Excel:"),
        excel_columns,
        ft.Divider(),
        generate_button,
        result_message
    )

# Ejecutar la aplicación
ft.app(target=main)